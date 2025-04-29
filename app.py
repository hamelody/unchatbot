import streamlit as st
import os
import fitz
import pandas as pd
import docx
from pptx import Presentation
import faiss
import openai
import numpy as np
import json
import io
from datetime import datetime
from openai import AzureOpenAI
from azure.storage.blob import BlobServiceClient

# --- 기본 세팅 ---
st.set_page_config(page_title="유앤생명과학 업무 가이드 봇", layout="centered")

# Azure OpenAI Client 초기화
openai_client = AzureOpenAI(
    api_key=st.secrets["AZURE_OPENAI_KEY"],
    azure_endpoint=st.secrets["AZURE_OPENAI_ENDPOINT"],
    api_version=st.secrets["AZURE_OPENAI_VERSION"]
)

# Azure Blob Storage 초기화
blob_service = BlobServiceClient.from_connection_string(st.secrets["AZURE_BLOB_CONN"])
container_client = blob_service.get_container_client(st.secrets["BLOB_CONTAINER"])

def save_to_blob(uploaded_file):
    data = uploaded_file.getvalue()
    blob_client = container_client.get_blob_client(uploaded_file.name)
    blob_client.upload_blob(data, overwrite=True)
    return f"https://{blob_service.account_name}.blob.core.windows.net/{container_client.container_name}/{uploaded_file.name}"

# 파일 경로 설정
INDEX_PATH = "vector.index"
LOG_PATH = "upload_log.json"
USAGE_LOG_PATH = "usage_log.json"
RULES_PATH = ".streamlit/prompt_rules.txt"

# 설정
EMBEDDING_MODEL = st.secrets["AZURE_OPENAI_EMBEDDING_DEPLOYMENT"]
ADMIN_PASSWORD = st.secrets["ADMIN_PASSWORD"]
TOKEN_COST = float(st.secrets.get("TOKEN_COST", 0.0))
CHUNK_SIZE = 500

# 벡터DB 초기화
if os.path.exists(INDEX_PATH):
    index = faiss.read_index(INDEX_PATH)
    with open("metadata.json", "r", encoding="utf-8") as f:
        metadata = json.load(f)
else:
    index = faiss.IndexFlatL2(1536)
    metadata = []

# 규칙 로드
def load_rules():
    if os.path.exists(RULES_PATH):
        return open(RULES_PATH, encoding="utf-8").read()
    return "당신은 제약회사 DI/GMP 전문가 챗봇입니다."

# 문서 파싱
def extract_text(file):
    ext = os.path.splitext(file.name)[1].lower()
    if ext == '.pdf':
        doc = fitz.open(stream=file.read(), filetype="pdf")
        return "\n".join(page.get_text() for page in doc)
    if ext == '.docx':
        doc_obj = docx.Document(file)
        return "\n".join(p.text for p in doc_obj.paragraphs)
    if ext in ('.xlsx', '.xlsm'):
        df = pd.read_excel(file)
        return df.to_string(index=False)
    if ext == '.csv':
        df = pd.read_csv(file)
        return df.to_string(index=False)
    if ext == '.pptx':
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    return ""

# 텍스트 청크 분할
def chunk_text(text):
    chunks, buf = [], ""
    for line in text.split("\n"):
        if len(buf) + len(line) < CHUNK_SIZE:
            buf += line + "\n"
        else:
            chunks.append(buf.strip())
            buf = line + "\n"
    if buf:
        chunks.append(buf.strip())
    return chunks

# 임베딩 호출
def get_embedding(text):
    resp = openai_client.embeddings.create(input=[text], model=EMBEDDING_MODEL)
    return resp.data[0].embedding

# 문서 DB 추가
def add_document(file, text, chunks):
    import numpy as np
    vecs, metas = [], []
    for c in chunks:
        emb = get_embedding(c)
        vecs.append(emb)
        metas.append({"file_name": file.name, "content": c})
    index.add(np.array(vecs).astype("float32"))
    metadata.extend(metas)
    faiss.write_index(index, INDEX_PATH)
    with open("metadata.json", "w", encoding="utf-8") as f:
        json.dump(metadata, f, ensure_ascii=False, indent=2)
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    log = {"file": file.name, "time": now, "chunks": len(chunks)}
    logs = json.load(open(LOG_PATH, encoding="utf-8")) if os.path.exists(LOG_PATH) else []
    logs.append(log)
    json.dump(logs, open(LOG_PATH, "w", encoding="utf-8"), ensure_ascii=False, indent=2)

# 문서 검색
def search_similar(query_or_chunk, k=5):
    qv = get_embedding(query_or_chunk)
    if index.ntotal == 0 or not metadata:
        return []
    D, I = index.search(np.array([qv]).astype("float32"), k)
    return [metadata[i]["content"] for i in I[0] if i < len(metadata)]

# --- UI 시작 ---
st.markdown("""
<div style="text-align:center; margin-bottom:24px;">
  <span style="font-size:2.1rem; font-weight:bold;">유앤생명과학 GMP/SOP 업무 가이드 봇</span>
  <span style="font-size:1rem; color:gray; margin-left:8px;">ver 0.5</span>
</div>
""", unsafe_allow_html=True)

# 탭 생성
tab2, tab1 = st.tabs(["💬 업무 질문", "⚙️ 관리자 설정"])

# --- 업무 질문 탭 ---
with tab2:
    st.header("업무 질문")
    st.markdown("""
    <div style='background-color:#edf4fb;padding:18px;border-radius:10px; margin-bottom:12px; max-width:700px; margin:auto;'>
      💡 예시: SOP 백업 주기, PIC/S Annex 11 차이 등
    </div>
    """, unsafe_allow_html=True)
    with st.form(key='query_form'):
        file_u = st.file_uploader("파일 첨부 (선택)", type=["pdf","docx","xlsx","xlsm","csv","pptx"])
        q = st.text_input("질문을 입력하세요...", key="input_text")
        submit = st.form_submit_button("전송")
    if submit and (q.strip() or file_u):
        with st.spinner("답변 생성 중..."):
            rules = load_rules()
            ctx = []
            if file_u:
                text = extract_text(file_u)
                chunks = chunk_text(text)
                for chunk in chunks:
                    ctx.extend(search_similar(chunk))
            else:
                ctx = search_similar(q)
            prompt = f"{rules}\n\n아래 문서를 참고해 답변하세요:\n" + "\n\n".join(ctx)
            resp = openai_client.chat.completions.create(
                model=st.secrets["AZURE_OPENAI_DEPLOYMENT"],
                messages=[{"role":"system","content":prompt},{"role":"user","content":q}],
                max_tokens=1000,
                temperature=0.2,
            )
            answer = resp.choices[0].message.content
            now = datetime.now().strftime("%Y-%m-%d %H:%M")
            st.markdown(f"""
            <div style='max-width:700px; margin:auto;'>
              <div style='text-align:right; margin-bottom:8px;'>
                <div style='display:inline-block; background-color:#f1f3f5; padding:10px; border-radius:15px; border-bottom-right-radius:0; max-width:80%;'>
                  <div style="font-size:12px; color:gray; text-align:right;">{now}</div>
                  {q}
                </div>
              </div>
              <div style='text-align:left;'>
                <div style='display:inline-block; background-color:#d4edda; padding:10px; border-radius:15px; border-bottom-left-radius:0; max-width:80%;'>
                  <div style="font-size:12px; color:gray; text-align:left;">{now}</div>
                  {answer}
                </div>
              </div>
            </div>
            """, unsafe_allow_html=True)
            st.session_state.last_q = q

# --- 관리자 설정 탭 ---
with tab1:
    st.header("⚙️ 관리자 설정")
    pw = st.text_input("관리자 비밀번호", type="password")
    if pw == ADMIN_PASSWORD:
        f = st.file_uploader("파일 업로드 (학습용)", type=["pdf","docx","xlsx","xlsm","csv","pptx"])
        if f:
            with st.spinner("처리 중..."):
                blob_url = save_to_blob(f)
                text = extract_text(f)
                chunks = chunk_text(text)
                add_document(f, text, chunks)
                st.success(f"✅ {f.name} 업로드 및 학습 완료 (Blob URL: {blob_url})")

        st.markdown("---")
        with st.expander("사용량 모니터링"):
            if os.path.exists(USAGE_LOG_PATH):
                use_logs = json.load(open(USAGE_LOG_PATH, encoding="utf-8"))
                df_use = pd.DataFrame(use_logs)
                total_calls = len(df_use)
                total_tokens = df_use["total_tokens"].sum()
                avg_tokens = total_tokens / total_calls if total_calls else 0
                total_cost = total_tokens * TOKEN_COST
                st.write(f"총 호출 횟수: {total_calls}")
                st.write(f"총 토큰 사용량: {total_tokens}")
                st.write(f"평균 토큰 사용량: {avg_tokens:.1f}")
                st.write(f"예상 비용 (USD): {total_cost:.4f}")
                st.dataframe(df_use)
            else:
                st.write("아직 호출 로그가 없습니다.")

        st.markdown("---")
        st.subheader("Blob에 저장된 파일 목록")
        blobs = container_client.list_blobs()
        blob_list = [{"파일명": b.name, "최종 수정": b.last_modified} for b in blobs]
        if blob_list:
            df_blob = pd.DataFrame(blob_list)
            st.dataframe(df_blob)
        else:
            st.write("Blob 컨테이너에 저장된 파일이 없습니다.")
    else:
        st.warning("관리자 비밀번호가 필요합니다.")